import os
from PyQt5 import QtWidgets
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from PIL import Image, UnidentifiedImageError
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
from modules.ui.view_dialog import InputDialog
from PyQt5.QtCore import Qt, QThread, pyqtSignal
from reportlab.pdfgen import canvas
from PyPDF2 import PdfWriter, PdfFileReader
import copy
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
import hashlib
from pathlib import Path
import pythoncom
import win32com.client
class WriteReport(QThread):
    progress_signal=pyqtSignal(int,str)
    result_signal = pyqtSignal()
    finished_signal=pyqtSignal()
    def __init__(self,date,name,duplicated_application):
        super().__init__()
        self.date=date
        self.name=name
        self.wiping_application=duplicated_application
        self.background_path='report_background.png'

    def ppt2pdf(self,ppt_target_file):
        pythoncom.CoInitialize()  # CoInitialize 호출 추가
        file_path = Path(ppt_target_file).resolve()
        out_file = file_path.parent / file_path.stem
        powerpoint = win32com.client.Dispatch("Powerpoint.Application")
        pdf = powerpoint.Presentations.Open(file_path, WithWindow=False)
        pdf.SaveAs(out_file, 32)
        pdf.Close()
        powerpoint.Quit()
        
    def copy_slide(self,sample,prs, index):
        template = sample.slides[index]
        blank_slide_layout = prs.slide_layouts[0]
        copied_slide = prs.slides.add_slide(blank_slide_layout)

        
        for shape in template.shapes:
            elem = shape.element
            new_elem = copy.deepcopy(elem)
            copied_slide.shapes._spTree.insert_element_before(new_elem, 'p:extLst')

        
        return copied_slide
    
    def select_shape_by_text(self,slide,text):
        for x in slide.shapes:
            if x.has_text_frame and x.text == text:
                return x
        print('요청한 Shape를 찾을 수 없습니다.')

    def set_background(self,slide, prs):
        """ 슬라이드의 배경으로 이미지를 설정하는 함수 """
        left = top = Inches(0)
        pic = slide.shapes.add_picture(self.background_path, left, top, width=prs.slide_width, height=prs.slide_height)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

    def calculate_file_hash(self,file_path, hash_algorithm='md5'):
        file_size = os.path.getsize(file_path)


        if hash_algorithm == 'md5':
            hasher = hashlib.md5()
        elif hash_algorithm == 'sha1':
            hasher = hashlib.sha1()
        elif hash_algorithm == 'sha256':
            hasher = hashlib.sha256()
        else:
            raise ValueError("지원되지 않는 해시 알고리즘입니다.")

        with open(file_path, 'rb') as file:
            while chunk := file.read(8192):
                hasher.update(chunk)

        hash_value = hasher.hexdigest()

        return file_size, hash_value



    def add_table_from_csv(self,sample,prs,slide,csvData, y_position,dataTitle):
        rows, cols = csvData.shape

        for start_row in range(0, rows, 10):
            end_row = min(start_row + 10, rows)
            
            table_width=6400800
            table_height = min(end_row - start_row, 10) * Inches(0.5)
            x_position = (prs.slide_width - table_width) / 2


            shape = slide.shapes.add_table(min(end_row-start_row, 10) + 1, cols, x_position, y_position, table_width, table_height)
            table=shape.table
            label=self.select_shape_by_text(slide,'title')
            label.text=dataTitle
            label.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

            for col_index, col_name in enumerate(csvData.columns):
                cell = table.cell(0, col_index)
                cell.text = col_name
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(8)

            for row_index in range(start_row, end_row):
                for col_index, item in enumerate(csvData.iloc[row_index]):
                    cell = table.cell(row_index - start_row + 1, col_index)
                    cell.text = str(item)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(8)

            tbl =  shape._element.graphic.graphicData.tbl
            style_id = '{616DA210-FB5B-4158-B5E0-FEB733F419BA}'

            tbl[0][-1].text = style_id
            if end_row!=rows:
                slide=self.copy_slide(sample,prs,2)
                self.set_background(slide, prs)
                y_position=Cm(5)

    def csvSlide(self,sample,prs,path,title):
        data=pd.read_csv(path)
        slide=self.copy_slide(sample,prs,2)
        self.set_background(slide, prs)
        y_position=Cm(5)
        self.add_table_from_csv(sample,prs,slide,data,y_position,title)

    def htmlSlide(self,sample,pre,path,dataTitle):
        fileList=os.listdir(path)
        for fileName in fileList:
            if '.txt' in fileName:
                filePath=os.path.join(path,fileName)
                with open(filePath, 'r', encoding='utf-8') as file:
                    text = file.read()
                    start = 0
                    max_lines = 35

                    while start < len(text):
                        slide=self.copy_slide(sample,pre,9)
                        self.set_background(slide, pre)
                        shape=self.select_shape_by_text(slide,'content')
                        title=self.select_shape_by_text(slide,'title')
                        title.text=dataTitle
                        label=self.select_shape_by_text(slide,'FileName')
                        label.text=fileName
                        label.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

                        content=''
                        for _ in range(max_lines):
                            if start < len(text):
                                end = text.find('\n', start) + 1 if text.find('\n', start) != -1 else len(text)
                                content += text[start:end]
                                start = end
                            else:
                                break
        
                        shape.text=content
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = RGBColor(0, 0, 0)  # 검은색
                                run.font.size = Pt(11)

    def ListSlide(self,sample,prs,path,dataTitle):
       
        readData=open(path,'r').read()
        dataList=readData.split('\n\n')
        for i in range(0,len(dataList),4):
            end_num=min(len(dataList)-i,4)
            slide=self.copy_slide(sample,prs,end_num+4)
            self.set_background(slide,prs)
            title=self.select_shape_by_text(slide,'title')
            title.text=dataTitle
            for image_num in range(end_num):
                shape=self.select_shape_by_text(slide,f'content{image_num+1}')
                label=self.select_shape_by_text(slide,f'ListName{image_num+1}')
                data_lines = dataList[i + image_num].split('\n')
                label.text=data_lines[0]
                shape.text='\n'.join(data_lines[1:])

                
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(0, 0, 0)  # 검은색
                        run.font.size = Pt(11)


                label.text_frame.paragraphs[0].runs[0].font.size=Pt(15)
                label.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    
    def ImageSlide(self,sample,prs,path,dataTitle):
        fileList=os.listdir(path)
        del fileList[-1]
        for i in range(0,len(fileList),2):
            end_num=min(len(fileList)-i,2)
            slide=self.copy_slide(sample,prs,end_num+2)
            self.set_background(slide,prs)
            title=self.select_shape_by_text(slide,'title')
            title.text=dataTitle
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            for image_num in range(end_num):
                shape=self.select_shape_by_text(slide,f'image{image_num+1}')
                label=self.select_shape_by_text(slide,f'data{image_num+1}')


                file_path=os.path.join(path,fileList[i+image_num])
                
                data=''
                data+=f"FILENAME: {os.path.basename(file_path)}\n"
                file_size, md5_hash = self.calculate_file_hash(file_path, 'md5')
                _, sha1_hash = self.calculate_file_hash(file_path, 'sha1')
                _, sha256_hash = self.calculate_file_hash(file_path, 'sha256')

                data+=f"FILESIZE: {file_size:,} bytes\n"
                data+=f"MD5: {md5_hash}\n"
                data+=f"SHA1: {sha1_hash}\n"
                data+=f"SHA256: {sha256_hash}\n"
                slide.shapes.add_picture(file_path,shape.left+Cm(0.5),shape.top+Cm(0.5),shape.width-Cm(1),shape.height-Cm(1))


                label.text=data
                for paragraph in label.text_frame.paragraphs:
                    paragraph.line_spacing = 1.5
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(0, 0, 0)  # 검은색
                        run.font.size = Pt(11)



    
    def UsageSlide(self,sample,prs):
        packages=pd.read_csv('./result/Package.csv')
        eventLog=pd.read_csv('./result/EventLog.csv')
        for i in range(len(packages)):
            row = packages.iloc[i]
            slide=self.copy_slide(sample,prs,1)
            packageShape=self.select_shape_by_text(slide,'Package')
            all_data_text = ''
            all_data_text += "\n".join([f"{column_name}: {value}" for column_name, value in row.items()])
            self.set_background(slide, prs)


            packageShape.text = all_data_text
            for paragraph in packageShape.text_frame.paragraphs:
                paragraph.line_spacing = 1.5
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(0, 0, 0)  # 검은색
                    run.font.size = Pt(11)

            eventlog_grouped = eventLog.groupby('new_group')
            eventlog_grouped=eventlog_grouped.get_group(i+1)


            y_position=packageShape.top+packageShape.height+Cm(1)
            self.add_table_from_csv(sample,prs,slide,eventlog_grouped,y_position,'UsageStats')
                
    
    def coverSlide(self,prs):
        slide_layout = prs.slide_layouts[5]
        cover_slide = prs.slides.add_slide(slide_layout)
        cover_slide.shapes.add_picture('./report_cover.png', 0, 0, prs.slide_width, prs.slide_height)
        original_title_y = Inches(2)
        title_shape = cover_slide.shapes.add_textbox(Inches(1), original_title_y + Inches(3), prs.slide_width - Inches(2), Inches(1))
        title_text_frame = title_shape.text_frame
        title_text_frame.text = "분석 결과 보고서"
        title_text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        title_paragraph = title_text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(44)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(0, 0, 100)
        original_subtitle_y = Inches(3)
        subtitle_shape = cover_slide.shapes.add_textbox(Inches(1), original_subtitle_y + Inches(3), prs.slide_width - Inches(2), Inches(1))
        subtitle_text_frame = subtitle_shape.text_frame
        subtitle_text_frame.text = f"Date: {self.date}   Name: {self.name}"
        subtitle_text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        return prs
    
    def run(self):
        signal=0

        sample = Presentation('sample.pptx')
        prs=Presentation()
        prs.slide_width = Cm(21)
        prs.slide_height = Cm(29.7)

        self.coverSlide(prs)

        signal+=10
        self.progress_signal.emit(signal,'UsageStats 보고서 작성 중')
        self.UsageSlide(sample,prs)
        signal+=10
        self.progress_signal.emit(signal,'Gallery Cache 보고서 작성 중')
        self.ImageSlide(sample,prs,'result/gallery3d_cache','Gallery Cache')
        signal+=10
        self.progress_signal.emit(signal,'Clipboard 보고서 작성 중')
        if os.path.exists('result/clipboard/clipboard.csv'):
            self.csvSlide(sample,prs,'result/clipboard/clipboard.csv','Clipboard')
            self.ImageSlide(sample,prs,'result/clipboard/image','Clipboard Image')
            self.htmlSlide(sample,prs,'result/clipboard/html','Clipboard HTML')
        elif os.path.exists('result/clipboard/clipboard.txt'):
            self.htmlSlide(sample,prs,'result/clipboard','Clipboard')

        if os.path.exists('result/contacts.csv'):
            signal+=10
            self.progress_signal.emit(signal,'Contacts 보고서 작성 중')
            self.csvSlide(sample,prs,'result/contacts.csv','Contacts')

        cacheList=os.listdir('result/cache')
        for cacheFolder in cacheList:
            signal+=10
            self.progress_signal.emit(signal,'Package Cache 보고서 작성 중')
            path=os.path.join('result/cache',cacheFolder)
            self.ImageSlide(sample,prs,path,cacheFolder)

        sharedPrefsList=os.listdir('result/shared_prefs')
        for fileName in sharedPrefsList:
            signal+=10
            self.progress_signal.emit(signal,'Shared Pref 보고서 작성 중')
            filePath=os.path.join('result/shared_prefs',fileName)
            self.ListSlide(sample,prs,filePath,fileName)
        
        prs.save('Analysis_Report.pptx')
        signal+=20
        self.progress_signal.emit(signal,'보고서 저장 중')
        # time.sleep(5)
        self.ppt2pdf('Analysis_Report.pptx')


        self.result_signal.emit()
        self.finished_signal.emit()


